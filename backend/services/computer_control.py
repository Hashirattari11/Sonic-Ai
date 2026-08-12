import ctypes
import os
import platform
import re
import subprocess
import time
from datetime import datetime

if platform.system() == "Windows":
    try:
        import win32gui
        _HAS_WIN32 = True
    except ImportError:
        _HAS_WIN32 = False
else:
    _HAS_WIN32 = False

# Virtual key codes (Windows media / volume keys)
VK_VOLUME_MUTE = 0xAD
VK_VOLUME_DOWN = 0xAE
VK_VOLUME_UP = 0xAF
VK_MEDIA_NEXT_TRACK = 0xB0
VK_MEDIA_PREV_TRACK = 0xB1
VK_MEDIA_STOP = 0xB2
VK_MEDIA_PLAY_PAUSE = 0xB3

ULONG_PTR = ctypes.c_ulonglong if ctypes.sizeof(ctypes.c_void_p) == 8 else ctypes.c_ulong


class KEYBDINPUT(ctypes.Structure):
    _fields_ = [("wVk", ctypes.c_ushort),
                ("wScan", ctypes.c_ushort),
                ("dwFlags", ctypes.c_ulong),
                ("time", ctypes.c_ulong),
                ("dwExtraInfo", ULONG_PTR)]


class MOUSEINPUT(ctypes.Structure):
    _fields_ = [("dx", ctypes.c_long),
                ("dy", ctypes.c_long),
                ("mouseData", ctypes.c_ulong),
                ("dwFlags", ctypes.c_ulong),
                ("time", ctypes.c_ulong),
                ("dwExtraInfo", ULONG_PTR)]


class HARDWAREINPUT(ctypes.Structure):
    _fields_ = [("uMsg", ctypes.c_ulong),
                ("wParamL", ctypes.c_ushort),
                ("wParamH", ctypes.c_ushort)]


class INPUTUNION(ctypes.Union):
    _fields_ = [("ki", KEYBDINPUT), ("mi", MOUSEINPUT), ("hi", HARDWAREINPUT)]


class INPUT(ctypes.Structure):
    _anonymous_ = ("u",)
    _fields_ = [("type", ctypes.c_ulong), ("u", INPUTUNION)]


def _send_key(vk_code: int):
    inputs = []
    for flags in (0, 2):  # press, then KEYEVENTF_KEYUP
        ki = KEYBDINPUT()
        ki.wVk = vk_code
        ki.wScan = 0
        ki.dwFlags = flags
        ki.time = 0
        ki.dwExtraInfo = 0
        inp = INPUT()
        inp.type = 1  # KEYBD_INPUT
        inp.ki = ki
        inputs.append(inp)
    ctypes.windll.user32.SendInput(2, (INPUT * 2)(*inputs), ctypes.sizeof(INPUT))


class SystemControl:
    """Handles low-level Windows/OS control: volume, media, brightness,
    windows, mouse, keyboard, power, files and more."""

    # ------------------------------------------------------------ volume
    @staticmethod
    def volume_up(steps: int = 2):
        try:
            for _ in range(steps):
                _send_key(VK_VOLUME_UP)
            return "Volume increased."
        except Exception as e:
            return f"Could not increase volume: {e}"

    @staticmethod
    def volume_down(steps: int = 2):
        try:
            for _ in range(steps):
                _send_key(VK_VOLUME_DOWN)
            return "Volume decreased."
        except Exception as e:
            return f"Could not decrease volume: {e}"

    @staticmethod
    def volume_mute():
        try:
            _send_key(VK_VOLUME_MUTE)
            return "Volume muted."
        except Exception as e:
            return f"Could not mute volume: {e}"

    @staticmethod
    def set_volume(level: int):
        """Set volume to an absolute level 0-100 using pycaw."""
        try:
            from pycaw.pycaw import AudioUtilities
            level = max(0, min(100, int(level)))
            volume = AudioUtilities.GetSpeakers().EndpointVolume
            volume.SetMasterVolumeLevelScalar(level / 100.0, None)
            return f"Volume set to {level}%."
        except Exception as e:
            return f"Could not set volume: {e}"

    @staticmethod
    def get_volume():
        try:
            from pycaw.pycaw import AudioUtilities
            volume = AudioUtilities.GetSpeakers().EndpointVolume
            level = round(volume.GetMasterVolumeLevelScalar() * 100)
            return {"level": level, "muted": bool(volume.GetMute())}
        except Exception:
            return None

    # ------------------------------------------------------- media keys
    @staticmethod
    def media_play_pause():
        try:
            _send_key(VK_MEDIA_PLAY_PAUSE)
            return "Media toggled play/pause."
        except Exception as e:
            return f"Could not toggle media: {e}"

    @staticmethod
    def media_next():
        try:
            _send_key(VK_MEDIA_NEXT_TRACK)
            return "Skipped to next track."
        except Exception as e:
            return f"Could not skip track: {e}"

    @staticmethod
    def media_prev():
        try:
            _send_key(VK_MEDIA_PREV_TRACK)
            return "Went to previous track."
        except Exception as e:
            return f"Could not go back: {e}"

    @staticmethod
    def media_stop():
        try:
            _send_key(VK_MEDIA_STOP)
            return "Media stopped."
        except Exception as e:
            return f"Could not stop media: {e}"

    # ------------------------------------------------------- brightness
    @staticmethod
    def set_brightness(level: int):
        try:
            import screen_brightness_control as sbc
            level = max(0, min(100, int(level)))
            sbc.set_brightness(level)
            return f"Brightness set to {level}%."
        except Exception as e:
            return f"Could not set brightness: {e}"

    @staticmethod
    def brightness_up(step: int = 10):
        try:
            import screen_brightness_control as sbc
            current = sbc.get_brightness()
            level = min(100, (current[0] if current else 0) + step)
            sbc.set_brightness(level)
            return f"Brightness set to {level}%."
        except Exception as e:
            return f"Could not change brightness: {e}"

    @staticmethod
    def brightness_down(step: int = 10):
        try:
            import screen_brightness_control as sbc
            current = sbc.get_brightness()
            level = max(0, (current[0] if current else 0) - step)
            sbc.set_brightness(level)
            return f"Brightness set to {level}%."
        except Exception as e:
            return f"Could not change brightness: {e}"

    # ---------------------------------------------------------- power
    @staticmethod
    def lock_screen():
        try:
            if platform.system() == "Windows":
                ctypes.windll.user32.LockWorkStation()
                return "Screen locked."
            return "Lock screen is only supported on Windows."
        except Exception as e:
            return f"Could not lock screen: {e}"

    @staticmethod
    def shutdown(delay: int = 0):
        try:
            if platform.system() == "Windows":
                subprocess.Popen(f"shutdown /s /t {int(delay)}", shell=True)
                return f"Computer will shut down in {delay} seconds."
            subprocess.Popen(["shutdown", "-h", "+1"])
            return "Shutting down..."
        except Exception as e:
            return f"Could not shut down: {e}"

    @staticmethod
    def restart(delay: int = 0):
        try:
            if platform.system() == "Windows":
                subprocess.Popen(f"shutdown /r /t {int(delay)}", shell=True)
                return f"Computer will restart in {delay} seconds."
            subprocess.Popen(["shutdown", "-r", "+1"])
            return "Restarting..."
        except Exception as e:
            return f"Could not restart: {e}"

    @staticmethod
    def sleep():
        try:
            if platform.system() == "Windows":
                subprocess.Popen("rundll32.exe powrprof.dll,SetSuspendState 0,1,0", shell=True)
                return "Putting computer to sleep..."
            return "Sleep is only supported on Windows."
        except Exception as e:
            return f"Could not sleep: {e}"

    @staticmethod
    def logoff():
        try:
            subprocess.Popen("shutdown /l", shell=True)
            return "Logging off..."
        except Exception as e:
            return f"Could not log off: {e}"

    @staticmethod
    def cancel_power():
        try:
            subprocess.Popen("shutdown /a", shell=True)
            return "Shutdown cancelled."
        except Exception as e:
            return f"Could not cancel shutdown: {e}"

    # ---------------------------------------------------------- apps
    @staticmethod
    def list_apps():
        try:
            result = subprocess.run(
                ["tasklist", "/FO", "CSV", "/NH"],
                capture_output=True, text=True, shell=True
            )
            names = {}
            for line in result.stdout.splitlines():
                parts = line.strip('"').split('","')
                if parts:
                    name = parts[0].strip('"').lower()
                    if name.endswith(".exe"):
                        names[name] = names.get(name, 0) + 1
            top = [n.replace(".exe", "") for n in sorted(names, key=lambda n: -names[n])[:20]]
            return {"handled": True, "response": "Running apps: " + ", ".join(top) if top else "No apps found."}
        except Exception as e:
            return {"handled": True, "response": f"Could not list apps: {e}"}

    @staticmethod
    def kill_app(name: str):
        try:
            name = name.lower().strip()
            if not name.endswith(".exe"):
                name += ".exe"
            result = subprocess.run(
                ["taskkill", "/IM", name, "/F"],
                capture_output=True, text=True, shell=True
            )
            if result.returncode == 0:
                return f"Closed {name.replace('.exe', '')}."
            return f"Could not close {name.replace('.exe', '')}. Is it running?"
        except Exception as e:
            return f"Could not close app: {e}"

    # ------------------------------------------------------- windows
    @staticmethod
    def _active_window():
        if not _HAS_WIN32:
            return None
        return win32gui.GetForegroundWindow()

    @staticmethod
    def minimize_window():
        try:
            if _HAS_WIN32:
                hwnd = win32gui.GetForegroundWindow()
                win32gui.ShowWindow(hwnd, 6)  # SW_MINIMIZE
                return "Window minimized."
            return "Window control is only supported on Windows."
        except Exception as e:
            return f"Could not minimize window: {e}"

    @staticmethod
    def maximize_window():
        try:
            if _HAS_WIN32:
                hwnd = win32gui.GetForegroundWindow()
                win32gui.ShowWindow(hwnd, 3)  # SW_MAXIMIZE
                return "Window maximized."
            return "Window control is only supported on Windows."
        except Exception as e:
            return f"Could not maximize window: {e}"

    @staticmethod
    def close_window():
        try:
            if _HAS_WIN32:
                hwnd = win32gui.GetForegroundWindow()
                win32gui.PostMessage(hwnd, 0x0010, 0, 0)  # WM_CLOSE
                return "Window closed."
            return "Window control is only supported on Windows."
        except Exception as e:
            return f"Could not close window: {e}"

    @staticmethod
    def show_desktop():
        try:
            import pyautogui
            pyautogui.hotkey("win", "d")
            return "Showing desktop."
        except Exception as e:
            return f"Could not show desktop: {e}"

    # --------------------------------------------------------- mouse
    @staticmethod
    def move_mouse(dx: int = 0, dy: int = 0, absolute=None):
        try:
            import pyautogui
            if absolute:
                pyautogui.moveTo(int(absolute[0]), int(absolute[1]), duration=0.3)
            else:
                pyautogui.moveRel(int(dx), int(dy), duration=0.3)
            return "Mouse moved."
        except Exception as e:
            return f"Could not move mouse: {e}"

    @staticmethod
    def click_mouse(button: str = "left"):
        try:
            import pyautogui
            pyautogui.click(button=button.lower())
            return f"{button} click performed."
        except Exception as e:
            return f"Could not click: {e}"

    @staticmethod
    def scroll(direction: str = "down"):
        try:
            import pyautogui
            amount = -300 if direction == "down" else 300
            pyautogui.scroll(amount)
            return f"Scrolled {direction}."
        except Exception as e:
            return f"Could not scroll: {e}"

    # ------------------------------------------------------- keyboard
    @staticmethod
    def type_text(text: str):
        try:
            import pyautogui
            pyautogui.write(text, interval=0.02)
            return f"Typed: {text[:60]}"
        except Exception as e:
            return f"Could not type: {e}"

    @staticmethod
    def press_hotkey(keys: str):
        try:
            import pyautogui
            key_list = [k.strip().lower() for k in keys.split("+") if k.strip()]
            pyautogui.hotkey(*key_list)
            return f"Pressed {keys}."
        except Exception as e:
            return f"Could not press keys: {e}"

    # --------------------------------------------------------- files
    @staticmethod
    def open_path(path: str):
        try:
            path = os.path.expandvars(os.path.expanduser(path))
            if not os.path.exists(path):
                return {"handled": True, "response": f"Path not found: {path}"}
            if platform.system() == "Windows":
                os.startfile(path)
            else:
                subprocess.Popen(["xdg-open", path])
            return {"handled": True, "response": f"Opened {path}."}
        except Exception as e:
            return {"handled": True, "response": f"Could not open path: {e}"}

    @staticmethod
    def common_folder(name: str):
        folders = {
            "downloads": "~\\Downloads",
            "download": "~\\Downloads",
            "desktop": "~\\Desktop",
            "documents": "~\\Documents",
            "document": "~\\Documents",
            "pictures": "~\\Pictures",
            "music": "~\\Music",
            "videos": "~\\Videos",
            "home": "~",
            "c": "C:\\",
            "c drive": "C:\\",
        }
        key = name.lower().strip()
        path = folders.get(key)
        if not path:
            return None
        return SystemControl.open_path(path)

    @staticmethod
    def create_folder(path: str):
        try:
            path = os.path.expandvars(os.path.expanduser(path))
            os.makedirs(path, exist_ok=True)
            return {"handled": True, "response": f"Created folder {path}."}
        except Exception as e:
            return {"handled": True, "response": f"Could not create folder: {e}"}

    @staticmethod
    def create_file(path: str, content: str = ""):
        try:
            path = os.path.expandvars(os.path.expanduser(path))
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            return {"handled": True, "response": f"Created file {path}."}
        except Exception as e:
            return {"handled": True, "response": f"Could not create file: {e}"}

    # ------------------------------------------------------- clipboard
    @staticmethod
    def get_clipboard():
        try:
            import pyperclip
            return pyperclip.paste()
        except Exception:
            try:
                result = subprocess.run(
                    ["powershell", "-Command", "Get-Clipboard"],
                    capture_output=True, text=True
                )
                return result.stdout.strip()
            except Exception:
                return None

    @staticmethod
    def set_clipboard(text: str):
        try:
            import pyperclip
            pyperclip.copy(text)
            return "Copied to clipboard."
        except Exception:
            try:
                encoded = text.replace("'", "''")
                subprocess.run(
                    ["powershell", "-Command", f"Set-Clipboard -Value '{encoded}'"],
                    capture_output=True
                )
                return "Copied to clipboard."
            except Exception as e:
                return f"Could not copy: {e}"

    # ------------------------------------------------------- screenshot
    @staticmethod
    def screenshot():
        try:
            import pyautogui
            img = pyautogui.screenshot()
            path = os.path.join(
                os.environ.get("TEMP", "/tmp"),
                f"sonic_screenshot_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png",
            )
            img.save(path)
            return {"handled": True, "response": f"Screenshot saved to {path}."}
        except Exception as e:
            return {"handled": True, "response": f"Screenshot failed: {e}"}


HINDI_HINTS = [
    "karo", "kar do", "kholo", "khol", "khol ke", "likho", "likh", "likh ke",
    "batao", "bata", "kya", "mujhe", "mera", "meri", "hai", "hain", "gaana",
    "gaane", "chalao", "dikhao", "dikha", "kaise", "kahan", "kab", "aur",
    "karke", "baare", "bare", "ke baare", "mein", "mai", "koi", "band",
    "shuru", "do", "de", "raha", "rahi", "so ja", "so jao", "utna",
]


def _is_hindi(text: str) -> bool:
    t = (text or "").lower()
    return any(h in t for h in HINDI_HINTS)


def _hi(en: str, ur: str) -> str:
    """Return the response in the language the user is likely speaking.
    Called with (english, hindi/urdu)."""
    return ur


class ComputerControl:
    """Voice-facing high-level command dispatcher.

    Returns a human-readable response string when a command was handled,
    or None when the utterance should be passed to the AI instead.
    Uses SystemControl under the hood (no nircmd required).
    """

    BROWSER_SITES = {
        "google": "https://www.google.com",
        "youtube": "https://www.youtube.com",
        "facebook": "https://www.facebook.com",
        "whatsapp": "https://web.whatsapp.com",
        "gmail": "https://mail.google.com",
        "github": "https://github.com",
        "chatgpt": "https://chatgpt.com",
        "instagram": "https://www.instagram.com",
        "twitter": "https://twitter.com",
        "linkedin": "https://www.linkedin.com",
        "netflix": "https://www.netflix.com",
        "spotify": "https://open.spotify.com",
        "amazon": "https://www.amazon.com",
        "stackoverflow": "https://stackoverflow.com",
        "maps": "https://www.google.com/maps",
        "weather": "https://www.google.com/search?q=weather",
    }

    APP_MAP = {
        "notepad": "notepad",
        "notepad app": "notepad",
        "calculator": "calc",
        "calc": "calc",
        "chrome": "chrome",
        "google chrome": "chrome",
        "browser": "chrome",
        "edge": "msedge",
        "microsoft edge": "msedge",
        "word": "winword",
        "microsoft word": "winword",
        "excel": "excel",
        "microsoft excel": "excel",
        "powerpoint": "powerpnt",
        "ppt": "powerpnt",
        "command prompt": "cmd",
        "cmd": "cmd",
        "terminal": "cmd",
        "paint": "mspaint",
        "task manager": "taskmgr",
        "file explorer": "explorer",
        "whatsapp": "whatsapp",
        "discord": "discord",
        "spotify": "spotify",
        "visual studio code": "code",
        "vs code": "code",
    }

    FILLER_WORDS = [
        "open", "youtube", "search", "for", "the", "song", "songs", "gaana",
        "gaane", "play", "karo", "kar", "kar do", "do", "per", "pe", "ka",
        "ki", "ke", "a", "an", "please", "plz", "mujhe", "koi", "ek", "some",
        "any",         "me", "mere", "meri", "batao", "bata", "chalao", "dhoondo",
        "khojo", "search kardo", "search karo", "kardo", "karna", "karo",
        "on", "ko", "se", "aur", "then", "and", "uske",
    ]

    def __init__(self):
        self._sys = SystemControl

    def handle(self, text: str):
        """Returns response string if handled, otherwise None."""
        text = (text or "").lower().strip()
        if not text:
            return None

        result = (
            self._handle_browser(text)
            or self._handle_youtube(text)
            or self._handle_song(text)
            or self._handle_websearch(text)
            or self._handle_open_app(text)
            or self._handle_volume(text)
            or self._handle_brightness(text)
            or self._handle_media(text)
            or self._handle_power(text)
            or self._handle_system(text)
        )
        return result

    # ------------------------------------------------------- helpers
    def _youtube_query(self, text):
        q = text
        for w in self.FILLER_WORDS:
            q = re.sub(rf"\b{re.escape(w)}\b", " ", q)
        q = re.sub(r"\s+", " ", q).strip(" ,-")
        return q

    def resolve_app(self, app: str):
        """Return a launchable target (exe path / command / url) for a
        friendly app name, or None. Tries App Paths registry, PATH, and
        common install locations."""
        import shutil
        name = (app or "").strip().lower()
        if not name:
            return None
        name = self.APP_MAP.get(name, name)

        # 1. Registry App Paths (covers chrome, msedge, office, etc.)
        import winreg
        exe = name if name.endswith(".exe") else f"{name}.exe"
        for hive, subkey in [
            (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths"),
            (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths"),
            (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths"),
        ]:
            try:
                with winreg.OpenKey(hive, subkey + "\\" + exe) as key:
                    path, _ = winreg.QueryValueEx(key, None)
                    if path and os.path.exists(path):
                        return path
            except OSError:
                pass

        # 2. PATH
        found = shutil.which(exe) or shutil.which(name)
        if found:
            return found

        # 3. Common install locations for known apps
        if name in self.APP_MAP:
            mapped = self.APP_MAP[name]
            mapped = mapped if mapped.endswith(".exe") else mapped + ".exe"
            candidates = []
            for base in [
                os.path.expandvars(r"%ProgramFiles%"),
                os.path.expandvars(r"%ProgramFiles(x86)%"),
                os.path.expandvars(r"%LOCALAPPDATA%"),
                os.path.expandvars(r"%APPDATA%"),
            ]:
                for rel in [
                    f"Google\\Chrome\\Application\\{mapped}",
                    f"Microsoft\\Edge\\Application\\{mapped}",
                    f"WhatsApp\\{mapped}",
                    f"Discord\\{mapped}",
                    f"Spotify\\{mapped}",
                    f"Microsoft Office\\root\\Office16\\{mapped}",
                ]:
                    p = os.path.join(base, rel)
                    candidates.append(p)
            # case-insensitive match against the candidate folder names
            for base in [
                os.path.expandvars(r"%ProgramFiles%"),
                os.path.expandvars(r"%ProgramFiles(x86)%"),
                os.path.expandvars(r"%LOCALAPPDATA%"),
                os.path.expandvars(r"%APPDATA%"),
            ]:
                if not os.path.isdir(base):
                    continue
                for folder in os.listdir(base):
                    if folder.lower() == name:
                        sub = os.path.join(base, folder)
                        for root, _dirs, files in os.walk(sub):
                            for fn in files:
                                if fn.lower() == mapped.lower():
                                    return os.path.join(root, fn)
                            break
            for p in candidates:
                if os.path.exists(p):
                    return p

        return None

    def open_app(self, app: str):
        """Launch an app by friendly name. Returns True on success (or when
        a web fallback was used for known web apps like whatsapp)."""
        try:
            import platform
            import webbrowser
            name = (app or "").strip().lower()
            if not name:
                return False
            if platform.system() != "Windows":
                subprocess.Popen(["open", "-a", name])
                return True
            exe = self.resolve_app(name)
            if exe:
                subprocess.Popen([exe], shell=False)
                return True
            # Web fallback for known web-app names (whatsapp, gmail, spotify web, etc.)
            web = self.BROWSER_SITES.get(name)
            if web:
                webbrowser.open(web)
                return True
            # fall back to 'start' (handles .lnk, aliases) then web version
            subprocess.Popen(["cmd", "/c", "start", "", name], shell=False)
            return True
        except Exception as e:
            print(f"open_app error: {e}")
            return False

    def type_text(self, text: str):
        return SystemControl.type_text(text)

    def _normalize_app(self, name: str):
        n = (name or "").strip().lower()
        n = n.replace("the ", "").replace("please", "").strip()
        if not n:
            return None
        if n in self.APP_MAP:
            return self.APP_MAP[n]
        if len(n.split()) == 1 and n.isalnum() or (len(n.split()) == 1 and "-" in n):
            return n
        return None

    def parse_write_task(self, text: str):
        """Detects 'open notepad and write about <topic>' / 'notepad khol ke X ke baare me likho'
        or 'make a presentation about X'. Returns {"app": ..., "topic": ...} or None."""
        t = (text or "").lower().strip()

        # Presentation task: "presentation bnao/make/create about X"
        is_presentation = re.search(r"(presentation|present|slide|powerpoint|ppt|pptx)", t)
        if is_presentation:
            has_make = any(w in t for w in ["bnao", "banao", "make", "create", "karo", "kar", "do", "bana", "build", "generate", "likho", "likh"])
            if has_make:
                return {"app": "presentation", "topic": self._extract_topic(t, "presentation")}

        if "notepad" not in t and "note pad" not in t:
            return None
        has_write = any(w in t for w in ["write", "likho", "likh", "type", "likh do", "type karo", "likh de"])
        if not has_write:
            return None

        topic = self._extract_topic(t, "notepad")
        if not topic:
            topic = "artificial intelligence"
        return {"app": "notepad", "topic": topic}

    def _extract_topic(self, t: str, mode: str) -> str:
        about_phrases = [
            "ke baare mein", "ke baare mai", "ke baare me", "ke bare mein",
            "ke bare mai", "ke bare me", "ke baare", "ke bare", "ka baare",
            "baare mein", "baare mai", "baare me", "baare", "about", "topic",
        ]
        topic = None
        for ap in about_phrases:
            if ap in t:
                if ap == "about" or ap == "topic":
                    topic = t.split(ap, 1)[1].strip()
                else:
                    before = t.split(ap, 1)[0].strip()
                    if mode == "notepad":
                        for kw in ["write", "likho", "likh de", "likh do", "likh", "type karo", "type"]:
                            if kw in before:
                                before = before.split(kw, 1)[1].strip()
                                break
                        before = re.sub(r"^(open\s+)?(the\s+)?(notepad|note pad)\s+(and\s+)?", "", before)
                        before = re.sub(r"\b(open|khol|kholo|khol ke|kholke|karke|karke|aur|and|then|the)\b", " ", before)
                    else:  # presentation
                        before = re.sub(r"^(open\s+)?(the\s+)?(make|create|build|bnao|banao|bana)?\s*(a\s+|an\s+)?(presentation|present|powerpoint|ppt)?\s*", "", before)
                        before = re.sub(r"\b(open|khol|kholo|make|create|build|bnao|banao|bana|presentation|present|powerpoint|ppt|aur|and|then|the|karo|kar)\b", " ", before)
                    before = re.sub(r"\s+", " ", before).strip(" ,")
                    topic = before
                break

        if not topic:
            kws = ["write", "likho", "likh de", "likh do", "likh", "type karo", "type", "make", "create", "bnao", "banao", "bana"]
            for kw in kws:
                if kw in t:
                    topic = t.split(kw, 1)[1].strip()
                    break

        topic = re.sub(r"^(open\s+)?(the\s+)?(notepad|note pad|presentation|powerpoint)\s+", "", (topic or ""))
        topic = re.sub(r"\b(open|khol|kholo|khol ke|kholke|karke|karo|kar do|make|create|build|bnao|banao|bana|presentation|present|powerpoint|ppt|aur|or|and|then|the|me|mai|mein|likh|likho|karo|do|ek|a|an)\b", " ", topic)
        topic = re.sub(r"\s+", " ", topic).strip()
        topic = re.sub(r"\s+(likh|likho|karo|kar|do|the|ka|ki|me|mai|mein)\s*$", "", topic).strip()
        topic = re.sub(r"\s+", " ", topic).strip()
        return topic

    def make_presentation(self, title: str, content_lines: list) -> str:
        """Build a real .pptx file on the Desktop and open it.
        Returns a human-readable status message."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            fname = f"SONIC_{title.replace(' ', '_')[:40]}.pptx"
            path = os.path.join(desktop, fname)

            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = f"{title} - SONIC AI"
            title_slide.placeholders[1].text = "Created by SONIC AI (J.A.R.V.I.S)"

            for i, line in enumerate(content_lines, start=1):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Slide {i}"
                lines = line.split("|")
                body = slide.placeholders[1].text_frame
                body.text = lines[0].strip() if lines else ""
                for extra in lines[1:]:
                    p = body.add_paragraph()
                    p.text = extra.strip()

            prs.save(path)

            import subprocess as _sp
            if os.name == "nt":
                os.startfile(path)
            else:
                _sp.Popen(["xdg-open", path])
            return path
        except Exception as e:
            return f"Could not create presentation: {e}"

    # ------------------------------------------------------- handlers
    def _handle_browser(self, text):
        if "open " not in text:
            return None
        remainder = text.split("open ", 1)[1].strip()
        remainder = remainder.replace("the ", "").strip()
        site = self.BROWSER_SITES.get(remainder)
        if not site:
            return None
        try:
            import webbrowser
            webbrowser.open(site)
            return _hi(f"Opening {remainder}.", f"{remainder.capitalize()} khol diya.")
        except Exception as e:
            return _hi(f"Could not open {remainder}: {e}", f"{remainder} nahi khul paya: {e}")

    def _handle_youtube(self, text):
        if "youtube" not in text and "you tube" not in text:
            return None
        query = self._youtube_query(text)
        try:
            import urllib.parse
            import webbrowser
            if query:
                url = f"https://www.youtube.com/results?search_query={urllib.parse.quote(query)}"
            else:
                url = "https://www.youtube.com"
            webbrowser.open(url)
            if query:
                return _hi(f"Searching YouTube for '{query}'.", f"YouTube pe '{query}' search kar raha hoon.")
            return _hi("Opening YouTube.", "YouTube khol raha hoon.")
        except Exception as e:
            return _hi(f"Could not open YouTube: {e}", f"YouTube nahi khul paya: {e}")

    def _handle_song(self, text):
        # media-control phrases must be handled by media handler, not YouTube search
        if any(k in text for k in [
            "next song", "next track", "skip song", "previous song",
            "previous track", "stop music", "stop song", "play music",
            "play pause", "play/pause", "gaana band",
        ]):
            return None
        is_song = any(k in text for k in ["song", "songs", "gaana", "gaane", "music"])
        is_search = any(k in text for k in ["search", "dhoondo", "khojo", "search kardo", "search karo", "batao", "play", "chalao", "youtube", "you tube"])
        if not is_song or not is_search:
            return None
        query = self._youtube_query(text)
        try:
            import urllib.parse
            import webbrowser
            if query:
                url = f"https://www.youtube.com/results?search_query={urllib.parse.quote(query)}"
            else:
                url = "https://www.youtube.com"
            webbrowser.open(url)
            if query:
                return _hi(f"Searching YouTube for '{query}'.", f"YouTube pe '{query}' search kar raha hoon.")
            return _hi("Opening YouTube.", "YouTube khol raha hoon.")
        except Exception as e:
            return _hi(f"Could not open YouTube: {e}", f"YouTube nahi khul paya: {e}")

    def _handle_websearch(self, text):
        if "search" not in text and "dhoondo" not in text and "google" not in text and "khojo" not in text:
            return None
        if any(k in text for k in ["song", "songs", "gaana", "gaane"]):
            return None
        q = text
        for w in ["search", "google", "for", "the", "web", "internet", "karo", "kar do", "kardo", "do", "pe", "per", "mein", "mai", "please", "plz", "kya", "mujhe", "me", "ko", "se", "batao", "bata", "jara", "baare", "bare"]:
            q = re.sub(rf"\b{re.escape(w)}\b", " ", q)
        q = re.sub(r"\s+", " ", q).strip(" ,-")
        try:
            import urllib.parse
            import webbrowser
            if not q:
                webbrowser.open("https://www.google.com")
                return _hi("Opening Google.", "Google khol raha hoon.")
            webbrowser.open(f"https://www.google.com/search?q={urllib.parse.quote(q)}")
            return _hi(f"Searching Google for '{q}'.", f"Google pe '{q}' search kar raha hoon.")
        except Exception as e:
            return _hi(f"Could not search: {e}", f"Search nahi ho paya: {e}")

    def _handle_open_app(self, text):
        open_words = ["open ", "launch ", "kholo ", "khol ", "khula ", "khol ke", "start "]
        matched = None
        for w in open_words:
            if w in text:
                matched = w
                break
        if matched is None:
            return None
        remainder = text.split(matched, 1)[1].strip()
        remainder = re.sub(r"\b(the|a|an|please|plz|karo|kar do|kardo|jara|mujhe)\b", " ", remainder).strip()
        if not remainder or remainder in ("this", "that", "the"):
            return None
        # Skip browser-site names handled above / known files
        if remainder in self.BROWSER_SITES:
            return None
        app = self._normalize_app(remainder)
        if app:
            if self.open_app(app):
                return _hi(f"Opening {remainder}.", f"{remainder.capitalize()} khol diya.")
            return _hi(f"Could not find {remainder} on your computer.", f"{remainder} aapke computer par nahi mila.")
        folder_res = self._sys.common_folder(remainder)
        if folder_res and isinstance(folder_res, dict) and folder_res.get("handled"):
            return folder_res["response"]
        return None

    def _handle_volume(self, text):
        try:
            if any(k in text for k in ["volume up", "increase volume", "louder", "awaz badhao", "aawaz badhao", "sound badhao"]):
                return self._sys.volume_up()
            if any(k in text for k in ["volume down", "decrease volume", "quieter", "awaz kam", "aawaz kam", "sound kam"]):
                return self._sys.volume_down()
            if "mute" in text or "awaz band" in text:
                return self._sys.volume_mute()
            if "unmute" in text or "awaz wapas" in text:
                return self._sys.volume_mute()
            if "volume" in text or "awaz" in text or "aawaz" in text:
                m = re.search(r"(\d{1,3})", text)
                if m and any(k in text for k in ["set", "karo", "kar", "raho", "pe", "per", "at", "to", "make"]):
                    return self._sys.set_volume(int(m.group(1)))
                if not any(k in text for k in ["up", "down", "mute", "unmute", "increase", "decrease", "loud", "quiet"]):
                    info = self._sys.get_volume()
                    if info:
                        state = "muted" if info["muted"] else f"at {info['level']}%"
                        return _hi(f"Volume is {state}.", f"Volume {state} hai.")
        except Exception as e:
            return _hi(f"Volume error: {e}", f"Volume set nahi ho paya: {e}")
        return None

    def _handle_brightness(self, text):
        try:
            if "brightness up" in text or "brightness badhao" in text:
                return self._sys.brightness_up()
            if "brightness down" in text or "brightness kam" in text:
                return self._sys.brightness_down()
            if "brightness" in text:
                m = re.search(r"(\d{1,3})", text)
                if m:
                    return self._sys.set_brightness(int(m.group(1)))
        except Exception as e:
            return _hi(f"Brightness error: {e}", f"Brightness set nahi ho paya: {e}")
        return None

    def _handle_media(self, text):
        try:
            if any(w in text for w in ["play music", "play pause", "play/pause", "gaana chalao", "song chalao"]):
                return self._sys.media_play_pause()
            if "next song" in text or "next track" in text or "skip song" in text or "aage wala gaana" in text:
                return self._sys.media_next()
            if "previous song" in text or "previous track" in text or "previous" in text or "piche wala gaana" in text:
                return self._sys.media_prev()
            if "stop music" in text or "stop song" in text or "gaana band" in text:
                return self._sys.media_stop()
            if text in ("play", "pause", "resume"):
                return self._sys.media_play_pause()
        except Exception as e:
            return _hi(f"Media error: {e}", f"Media control nahi ho paya: {e}")
        return None

    def _handle_power(self, text):
        if "lock" in text and "screen" in text or "screen lock" in text:
            return self._sys.lock_screen()
        if "shut down" in text or "shutdown" in text or "band karo computer" in text:
            return self._sys.shutdown(30)
        if "restart" in text or "reboot" in text:
            return self._sys.restart(30)
        if "sleep" in text or "sula do" in text:
            return self._sys.sleep()
        return None

    def _handle_system(self, text):
        if "screenshot" in text or "capture screen" in text or "screenshot le lo" in text:
            res = self._sys.screenshot()
            return res["response"]
        if "what time is it" in text or "current time" in text or "what's the time" in text or "time batao" in text or "kitne baje" in text:
            now = datetime.now()
            return _hi(f"The time is {now.strftime('%I:%M %p')}.", f"Abhi time {now.strftime('%I:%M')} ho raha hai.")
        if "what is the date" in text or "today's date" in text or "what date is it" in text or "aaj ki date" in text or "date batao" in text:
            return _hi(f"Today is {datetime.now().strftime('%A, %B %d, %Y')}.", f"Aaj {datetime.now().strftime('%d %B %Y')} hai.")
        return None
